"""
Document Processing Service
Handles text extraction, OCR, cleaning, chunking, and metadata extraction
"""
import os
import re
from typing import List, Dict, Optional, Tuple
from pathlib import Path
import PyPDF2
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image
import pytesseract
from io import BytesIO


class DocumentProcessor:
    """Process documents for RAG pipeline"""
    
    def __init__(self):
        self.supported_formats = {
            'pdf': self._extract_pdf,
            'docx': self._extract_docx,
            'pptx': self._extract_pptx,
            'xlsx': self._extract_xlsx,
            'txt': self._extract_txt,
            'md': self._extract_txt,
        }
    
    def extract_text(self, file_path: str, file_type: str) -> str:
        """Extract text from document based on file type"""
        file_type = file_type.lower().lstrip('.')
        
        if file_type in self.supported_formats:
            return self.supported_formats[file_type](file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    
    def extract_text_from_image(self, image_path: str) -> str:
        """Extract text from image using OCR"""
        try:
            image = Image.open(image_path)
            text = pytesseract.image_to_string(image)
            return text
        except Exception as e:
            raise Exception(f"OCR failed: {str(e)}")
    
    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            raise Exception(f"PDF extraction failed: {str(e)}")
        return text
    
    def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX"""
        text = ""
        try:
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
        except Exception as e:
            raise Exception(f"DOCX extraction failed: {str(e)}")
        return text
    
    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PPTX"""
        text = ""
        try:
            ppt = Presentation(file_path)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            raise Exception(f"PPTX extraction failed: {str(e)}")
        return text
    
    def _extract_xlsx(self, file_path: str) -> str:
        """Extract text from XLSX"""
        text = ""
        try:
            wb = load_workbook(file_path)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) if cell is not None else "" for cell in row])
                    text += row_text + "\n"
        except Exception as e:
            raise Exception(f"XLSX extraction failed: {str(e)}")
        return text
    
    def _extract_txt(self, file_path: str) -> str:
        """Extract text from TXT or MD"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read()
            except Exception as e:
                raise Exception(f"Text extraction failed: {str(e)}")
    
    def clean_text(self, text: str) -> str:
        """Clean and normalize extracted text"""
        # Remove extra whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove special characters but keep basic punctuation
        text = re.sub(r'[^\w\s\.\,\!\?\;\:\-\(\)]', '', text)
        
        # Remove multiple consecutive punctuation
        text = re.sub(r'([\.!?])\1+', r'\1', text)
        
        # Remove page numbers and headers (common patterns)
        text = re.sub(r'\n\s*\d+\s*\n', '\n', text)
        text = re.sub(r'Page \d+ of \d+', '', text)
        
        # Strip leading/trailing whitespace
        text = text.strip()
        
        return text
    
    def chunk_text(self, text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
        """Split text into overlapping chunks"""
        if not text:
            return []
        
        chunks = []
        start = 0
        text_length = len(text)
        
        while start < text_length:
            end = start + chunk_size
            
            # If this is the last chunk, take whatever is left
            if end >= text_length:
                chunks.append(text[start:])
                break
            
            # Try to break at a sentence boundary
            chunk = text[start:end]
            last_period = chunk.rfind('.')
            last_question = chunk.rfind('?')
            last_exclamation = chunk.rfind('!')
            
            # Find the last sentence ending
            last_sentence_end = max(last_period, last_question, last_exclamation)
            
            if last_sentence_end > chunk_size * 0.7:  # Only break if it's not too early
                end = start + last_sentence_end + 1
                chunk = text[start:end]
            
            chunks.append(chunk)
            start = end - overlap
        
        return chunks
    
    def extract_metadata(self, file_path: str, file_type: str) -> Dict:
        """Extract metadata from document"""
        metadata = {
            'file_name': os.path.basename(file_path),
            'file_type': file_type,
            'file_size': os.path.getsize(file_path),
            'page_count': None,
            'word_count': None,
            'char_count': None,
        }
        
        try:
            text = self.extract_text(file_path, file_type)
            metadata['word_count'] = len(text.split())
            metadata['char_count'] = len(text)
            
            if file_type.lower() == 'pdf':
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    metadata['page_count'] = len(pdf_reader.pages)
        except Exception:
            pass
        
        return metadata
    
    def process_document(self, file_path: str, file_type: str) -> Dict:
        """Complete document processing pipeline"""
        # Extract text
        text = self.extract_text(file_path, file_type)
        
        # Clean text
        cleaned_text = self.clean_text(text)
        
        # Extract metadata
        metadata = self.extract_metadata(file_path, file_type)
        
        # Chunk text
        chunks = self.chunk_text(cleaned_text)
        
        return {
            'text': cleaned_text,
            'chunks': chunks,
            'metadata': metadata,
            'chunk_count': len(chunks)
        }
